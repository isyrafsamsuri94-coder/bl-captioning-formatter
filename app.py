import re
from dataclasses import dataclass
from typing import List, Set, Optional

import streamlit as st
import spacy
from spacy.matcher import PhraseMatcher
from pptx import Presentation
import io

# ============================================================
# spaCy model loading (cached for Streamlit Cloud)
# ============================================================

@st.cache_resource
def load_nlp():
    return spacy.load("en_core_web_sm")

nlp = load_nlp()

# ============================================================
# Configuration
# ============================================================

@dataclass
class SubtitleCfg:
    max_len: int = 120
    max_lines: int = 1
    block_budget: int = 120
    preserve_paragraph_breaks: bool = True

# ============================================================
# Normalization
# ============================================================

def normalize_ws(text: str) -> str:
    """Normalize whitespace only; preserve punctuation and words."""
    return re.sub(r"[ \t]+", " ", text).strip()

# ============================================================
# Phrase protection configuration
# ============================================================

NEVER_BREAK_PHRASES = [
    "as well as",
    "in order to",
    "due to",
    "because of",
    "as a result",
    "in terms of",
    "according to",
    "rather than",
    "even though",
    "with respect to",
    "in relation to",
    "based on",
    "at least",
    "more than",
    "less than",
    "no more than",
    "as soon as possible",
    "at the end of the day",
    "in the meantime",
    "for the most part",
    "more or less",
    "goes through",
]

TITLES = {"mr", "mrs", "ms", "dr", "prof", "professor", "sir", "madam"}

UNITS = {
    "percent", "%", "km", "m", "cm", "mm", "kg", "g", "mg", "hz",
    "seconds", "second", "mins", "min", "minutes", "minute",
    "hours", "hour", "dollars", "dollar", "usd", "sgd", "eur", "gbp"
}

# ============================================================
# Protected spans
# ============================================================

def protected_spans(doc) -> List:
    spans = []

    # Phrase matcher
    matcher = PhraseMatcher(nlp.vocab, attr="LOWER")
    patterns = [nlp.make_doc(p) for p in NEVER_BREAK_PHRASES]
    matcher.add("NBP", patterns)

    for _, start, end in matcher(doc):
        spans.append(doc[start:end])

    # Named entities
    spans.extend(list(doc.ents))

    # Proper noun sequences
    i = 0
    while i < len(doc):
        if i + 1 < len(doc) and doc[i].pos_ == "PROPN" and doc[i+1].pos_ == "PROPN":
            start = i
            while i < len(doc) and doc[i].pos_ == "PROPN":
                i += 1
            spans.append(doc[start:i])
        else:
            i += 1

    return spans

def forbid_breaks_inside_spans(spans) -> Set[int]:
    forbid = set()
    for sp in spans:
        for tok in sp[:-1]:
            forbid.add(tok.i)
    return forbid

# ============================================================
# Break prohibition rules
# ============================================================

def forbidden_break_token_indices(doc) -> Set[int]:
    forbid: Set[int] = set()

    spans = protected_spans(doc)
    forbid |= forbid_breaks_inside_spans(spans)

    for i, tok in enumerate(doc[:-1]):
        # Compound nouns
        if tok.dep_ == "compound":
            forbid.add(i)

        # Verb cohesion
        if tok.dep_ in {"aux", "neg", "prt"}:
            forbid.add(i)

        # Prepositions: always forbid breaking after ("prep | rest").
        # Also forbid breaking before ("word | prep rest") only when the
        # preposition is directly governed by the immediately preceding token
        # (e.g. "distributed with", "coincide with", "member of").
        # This avoids blocking clause-level breaks before prepositions like
        # "to" in "from X to Y", where the head is far away.
        if tok.pos_ == "ADP":
            forbid.add(i)
            if i > 0 and tok.head.i == i - 1:
                forbid.add(i - 1)

        # Titles and units
        if tok.lower_ in TITLES or tok.lower_ in UNITS:
            forbid.add(i)

        # 's-genitives: never split "Cas9's | nuclease" or "Cas9 | 's"
        if re.search(r"[\u2019']s$", tok.text) and i > 0:
            forbid.add(i - 1)  # before 's
            forbid.add(i)      # after 's

    return forbid

# ============================================================
# Helpers
# ============================================================

def is_break_after_function_word(line1: str) -> bool:
    return bool(
        re.search(
            r"\b(the|a|an|of|to|in|on|at|for|as|by|with|and|or|but)\s*$",
            line1.lower(),
        )
    )

def punctuation_bonus(line1: str) -> int:
    if re.search(r"[.!?]\s*$", line1):
        return 12
    if re.search(r"[;:]\s*$", line1):
        return 8
    if re.search(r",\s*$", line1):
        return 4
    return 0

# ============================================================
# Transcript segmentation
# ============================================================

def split_into_paragraphs(text: str, preserve: bool) -> List[str]:
    text = text.strip()
    if not preserve:
        return [normalize_ws(text)]
    return [
        p.strip()
        for p in re.split(r"\n\s*\n", text)
        if p.strip()
    ]

def _regex_split_sentences(text: str) -> List[str]:
    """Split on obvious sentence boundaries that spaCy may miss."""
    parts = re.split(r'(?<=[.!?])\s+(?=[A-Z\u201c"])', text)
    return [p.strip() for p in parts if p.strip()]

def sentence_units(paragraph: str) -> List[str]:
    para = normalize_ws(paragraph)
    if not para:
        return []
    doc = nlp(para)
    spacy_sents = [normalize_ws(sent.text) for sent in doc.sents if sent.text.strip()]
    result = []
    for s in spacy_sents:
        result.extend(_regex_split_sentences(s))
    return result

def clause_splits(sentence: str) -> List[str]:
    if len(sentence) > 180:
        parts = re.split(r"[,;:]", sentence)
        return [normalize_ws(p) for p in parts if p.strip()]
    return [sentence]

def pack_units_into_blocks(units: List[str], cfg: SubtitleCfg) -> List[str]:
    blocks = []
    buf = ""

    for u in units:
        candidate = f"{buf} {u}".strip() if buf else u
        if len(candidate) <= cfg.block_budget:
            buf = candidate
        else:
            if buf:
                blocks.append(buf)
            buf = u

    if buf:
        blocks.append(buf)

    return blocks

# ============================================================
# Smart recursive block splitter
# ============================================================

def split_block_to_lines(text: str, max_len: int) -> List[str]:
    """Recursively split text into lines each <= max_len at appropriate boundaries."""
    text = normalize_ws(text)
    if len(text) <= max_len:
        return [text]

    doc = nlp(text)
    forbid = forbidden_break_token_indices(doc)

    best_score = -float('inf')
    best_left: Optional[str] = None
    best_right: Optional[str] = None

    for i, tok in enumerate(doc[:-1]):
        if tok.i in forbid:
            continue

        split_pos = tok.idx + len(tok.text)
        left = text[:split_pos].rstrip()
        right = text[split_pos:].lstrip()

        if not right or len(right) < 10:
            continue
        if is_break_after_function_word(left):
            continue

        # Prefer punctuation boundaries; penalise very unbalanced splits
        score = punctuation_bonus(left) * 100 - abs(len(left) - len(right))
        if score > best_score:
            best_score = score
            best_left = left
            best_right = right

    if best_left is None:
        # Hard fallback: split at word boundary near max_len
        idx = text.rfind(" ", 0, max_len)
        if idx == -1:
            idx = max_len
        best_left = text[:idx].strip()
        best_right = text[idx:].strip()

    result = split_block_to_lines(best_left, max_len)
    if best_right:
        result.extend(split_block_to_lines(best_right, max_len))
    return result

# ============================================================
# Merge orphaned conjunctions
# ============================================================

_ORPHAN_CONJUNCTIONS = {"and", "or", "but", "nor", "yet", "so"}

def merge_orphaned_conjunctions(blocks: List[str], cfg: SubtitleCfg) -> List[str]:
    """Join a lone conjunction block with the following content block."""
    result: List[str] = []
    i = 0
    while i < len(blocks):
        block = blocks[i]
        if block != "" and block.strip().lower() in _ORPHAN_CONJUNCTIONS:
            # Find the next non-empty block
            j = i + 1
            while j < len(blocks) and blocks[j] == "":
                j += 1
            if j < len(blocks):
                merged = f"{block} {blocks[j]}"
                if len(merged) <= cfg.block_budget:
                    result.append(merged)
                    i = j + 1
                    continue
        result.append(block)
        i += 1
    return result

# ============================================================
# Remove trailing commas
# ============================================================

def remove_trailing_commas_from_blocks(blocks: List[str]) -> List[str]:
    out = []
    for b in blocks:
        if b == "":
            out.append(b)
            continue
        lines = b.split("\n")
        lines = [re.sub(r",\s*$", "", ln) for ln in lines]
        out.append("\n".join(lines))
    return out

# ============================================================
# PPTX notes extraction
# ============================================================
def extract_notes(pptx_file):
    prs = Presentation(pptx_file)
    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
            if text:
                notes.append(text)
    return notes
# ============================================================
# Main formatter
# ============================================================

def format_transcript_rules(text: str, cfg: SubtitleCfg) -> List[str]:
    blocks_out = []

    paragraphs = split_into_paragraphs(text, cfg.preserve_paragraph_breaks)

    for p in paragraphs:
        sentences = sentence_units(p)
        for sent in sentences:
            clauses = clause_splits(sent)
            blocks = pack_units_into_blocks(clauses, cfg)
            for b in blocks:
                if len(b) > cfg.max_len:
                    blocks_out.extend(split_block_to_lines(b, cfg.max_len))
                else:
                    blocks_out.append(b)
            blocks_out.append("")

    if blocks_out and blocks_out[-1] == "":
        blocks_out.pop()

    return blocks_out

def format_transcript_hybrid(text: str, cfg: SubtitleCfg) -> str:
    blocks = format_transcript_rules(text, cfg)
    blocks = remove_trailing_commas_from_blocks(blocks)
    blocks = merge_orphaned_conjunctions(blocks, cfg)
    # Join non-empty blocks with exactly one blank line between them
    content_blocks = [b for b in blocks if b != ""]
    return "\n\n".join(content_blocks)

# ============================================================
# Streamlit UI
# ============================================================

st.title("BL Captioning Formatter (Beta)")
st.markdown(
    """
    Upload `.txt` file to generate formatted captions for Descript 
    \nOR 
    \nUpload `.pptx` file to extract speaker notes and optionally format them into captions.

    > ⚠️ **You must review the output carefully before using it as captions.**
    > The formatting may not be perfect and could break words or split sentences inappropriately in some cases.
    > This tool is intended to save time on manual caption formatting, but cannot replace human judgment.
    """
)

mode = st.selectbox(
    "What would you like to do?",
    ["Select an option...", "Format captions from TXT", "Extract speaker notes from PPTX"],
)

if mode == "Format captions from TXT":
    uploaded_file = st.file_uploader("Upload a .txt file", type=["txt"])
elif mode == "Extract speaker notes from PPTX":
    uploaded_file = st.file_uploader("Upload a .pptx file", type=["pptx"])
else:
    uploaded_file = None

if uploaded_file is not None:

    cfg = SubtitleCfg(
    max_len=110,
    max_lines=2,
    block_budget=110,
    preserve_paragraph_breaks=False,
    )

    if mode == "Extract speaker notes from PPTX":
        slides_text = extract_notes(io.BytesIO(uploaded_file.read()))
        output_text = "\n\n".join(slides_text)

        st.success("Extraction complete!")
        st.subheader("Preview")
        st.text_area("", value=output_text, height=400, disabled=True, label_visibility="collapsed")

        st.download_button(
            label="Download speaker notes",
            data=output_text,
            file_name="speaker_notes.txt",
            mime="text/plain",
        )

        formatextractednotes = st.selectbox(
        "Would you like to format the extracted notes into captions?",
        ["No", "Yes"]
        )

        if formatextractednotes == "Yes":
            input_text = output_text

            with st.spinner("Formatting captions..."):
                output_text = format_transcript_hybrid(input_text, cfg)

            st.success("Formatting complete!")
            st.subheader("Preview")
            st.text_area("", value=output_text, height=400, disabled=True, label_visibility="collapsed")

            st.download_button(
                label="Download formatted captions",
                data=output_text,
                file_name="formatted_captions.txt",
                mime="text/plain",
            )

    else:
        input_text = uploaded_file.read().decode("utf-8")

        with st.spinner("Formatting captions..."):
            output_text = format_transcript_hybrid(input_text, cfg)

        st.success("Formatting complete!")

        st.subheader("Preview")
        st.text_area("", value=output_text, height=400, disabled=True, label_visibility="collapsed")

        st.download_button(
            label="Download formatted text",
            data=output_text,
            file_name="formatted_captions.txt",
            mime="text/plain",
        )